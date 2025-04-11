import os
os.environ["TOKENIZERS_PARALLELISM"] = "false"

import streamlit as st
import json
import faiss
import numpy as np
import torch
import subprocess
import fitz  # PyMuPDF for PDF parsing
from io import BytesIO
from sentence_transformers import SentenceTransformer

import docx
from bs4 import BeautifulSoup
from pptx import Presentation
import pytesseract
from PIL import Image

#########################
# Utility Functions
#########################
def get_device():
    """Return 'mps' if available on macOS, else 'cpu'."""
    if torch.backends.mps.is_available():
        return "mps"
    else:
        return "cpu"

#########################
# File Processing Functions for Multiple Formats
#########################

def process_pdf(file_bytes, heading_threshold=14):
    """Process an uploaded PDF file using PyMuPDF."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages_data = []
    for page in doc:
        page_number = page.number + 1  # 1-indexed
        page_dict = page.get_text("dict")
        page_text_lines = []
        headings = []
        for block in page_dict.get("blocks", []):
            if "lines" not in block:
                continue
            block_text = ""
            for line in block["lines"]:
                line_text = ""
                for span in line["spans"]:
                    text = span["text"].strip()
                    if not text:
                        continue
                    line_text += text + " "
                    if span["size"] >= heading_threshold and text not in headings:
                        headings.append(text)
                block_text += line_text.strip() + "\n"
            page_text_lines.append(block_text.strip())
        pages_data.append({
            "page_number": page_number,
            "headings": headings,
            "text": "\n".join(page_text_lines)
        })
    return pages_data

def process_text(file_bytes):
    """Process a plain text file."""
    text = file_bytes.decode("utf-8")
    return [{
        "page_number": 1,
        "headings": [],
        "text": text
    }]

def process_docx(file_bytes):
    """Process a DOCX file using python-docx."""
    document = docx.Document(BytesIO(file_bytes))
    text = "\n".join([para.text for para in document.paragraphs])
    return [{
        "page_number": 1,
        "headings": [],
        "text": text
    }]

def process_html(file_bytes):
    """Process an HTML file using BeautifulSoup."""
    soup = BeautifulSoup(file_bytes, "html.parser")
    text = soup.get_text(separator="\n")
    return [{
        "page_number": 1,
        "headings": [],
        "text": text
    }]

def process_pptx(file_bytes):
    """Process a PPTX file using python-pptx (each slide as a 'page')."""
    prs = Presentation(BytesIO(file_bytes))
    pages_data = []
    slide_num = 1
    for slide in prs.slides:
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
        slide_text = "\n".join(text_runs)
        pages_data.append({
            "page_number": slide_num,
            "headings": [],
            "text": slide_text
        })
        slide_num += 1
    return pages_data

def process_image(file_bytes):
    """Process an image file via OCR using pytesseract."""
    img = Image.open(BytesIO(file_bytes))
    text = pytesseract.image_to_string(img)
    return [{
        "page_number": 1,
        "headings": [],
        "text": text
    }]

def process_file(uploaded_file):
    """
    Dispatch processing based on file extension.
    Supported extensions: pdf, txt, docx, html, pptx, png, jpg, jpeg, tiff.
    """
    filename = uploaded_file.name
    ext = filename.split(".")[-1].lower()
    file_bytes = uploaded_file.read()
    if ext == "pdf":
        pages_data = process_pdf(file_bytes)
    elif ext == "txt":
        pages_data = process_text(file_bytes)
    elif ext == "docx":
        pages_data = process_docx(file_bytes)
    elif ext in ["html", "htm"]:
        pages_data = process_html(file_bytes)
    elif ext == "pptx":
        pages_data = process_pptx(file_bytes)
    elif ext in ["png", "jpg", "jpeg", "tiff"]:
        pages_data = process_image(file_bytes)
    else:
        st.warning(f"Unsupported file type: {ext}")
        pages_data = []
    # Add document name metadata to each page.
    for page in pages_data:
        page["document_name"] = filename
    return pages_data

#########################
# Text Segmentation and Chunking
#########################
def segment_page_text(page_data):
    """
    Split a page's text using detected headings as markers.
    If no headings are found, treat the whole page as one chunk.
    """
    page_number = page_data["page_number"]
    text = page_data["text"]
    headings = page_data.get("headings", [])
    chunks = []
    if headings:
        lines = text.splitlines()
        current_section = None
        current_chunk_lines = []
        for line in lines:
            if line.strip() in headings:
                if current_section is not None and current_chunk_lines:
                    chunks.append({
                        "document_name": page_data.get("document_name", "Unknown"),
                        "page_number": page_number,
                        "section_title": current_section,
                        "chunk_text": "\n".join(current_chunk_lines).strip()
                    })
                current_section = line.strip()
                current_chunk_lines = []
            else:
                current_chunk_lines.append(line)
        if current_section is not None and current_chunk_lines:
            chunks.append({
                "document_name": page_data.get("document_name", "Unknown"),
                "page_number": page_number,
                "section_title": current_section,
                "chunk_text": "\n".join(current_chunk_lines).strip()
            })
        if not chunks:
            chunks.append({
                "document_name": page_data.get("document_name", "Unknown"),
                "page_number": page_number,
                "section_title": headings[0] if headings else f"Page {page_number}",
                "chunk_text": text
            })
    else:
        chunks.append({
            "document_name": page_data.get("document_name", "Unknown"),
            "page_number": page_number,
            "section_title": f"Page {page_number}",
            "chunk_text": text
        })
    return chunks

def segment_document(pages_data):
    """Segment all pages into chunks with metadata."""
    all_chunks = []
    for page in pages_data:
        chunks = segment_page_text(page)
        all_chunks.extend(chunks)
    return all_chunks

#########################
# Embedding and Index Construction
#########################
def compute_embeddings(chunks, model_name='sentence-transformers/all-MiniLM-L6-v2'):
    """Compute embeddings for each chunk using the MiniLM model."""
    device = get_device()
    model = SentenceTransformer(model_name, device=device)
    documents = [chunk['chunk_text'] for chunk in chunks]
    embeddings = model.encode(documents, convert_to_numpy=True)
    return embeddings

def build_faiss_index(embeddings):
    """Build a FAISS index from the embeddings."""
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings)
    return index

#########################
# Query Processing Functions
#########################
def compute_query_embedding(query, model_name='sentence-transformers/all-MiniLM-L6-v2'):
    device = get_device()
    model = SentenceTransformer(model_name, device=device)
    return model.encode([query], convert_to_numpy=True)

def retrieve_relevant_chunks(query, index, metadata, top_k=5):
    """Retrieve the top_k most relevant chunks for the query."""
    query_embedding = compute_query_embedding(query)
    distances, indices = index.search(query_embedding, top_k)
    relevant_chunks = []
    for idx in indices[0]:
        if idx != -1 and idx < len(metadata):
            relevant_chunks.append(metadata[idx])
    return relevant_chunks

def construct_prompt(query, retrieved_chunks, max_chunk_length=500, conversation_context=""):
    """
    Build a prompt that includes:
     - (Optional) conversation context (previous Q&A)
     - The current user query
     - Instructions to skip chain-of-thought and return only the final answer
     - Relevant excerpts (truncated) from the uploaded documents.
    """
    prompt = ""
    if conversation_context:
        prompt += conversation_context + "\n\n"
    prompt += f"User Question: {query}\n\n"
    prompt += ("Below are relevant excerpts from the uploaded documents. "
               "Answer the question based solely on the above context. "
               "Do not show any chain-of-thought or thinking process – only provide your final answer.\n")
    for chunk in retrieved_chunks:
        truncated_text = chunk['chunk_text'][:max_chunk_length]
        prompt += (
            f"\n[Document: {chunk['document_name']}, Page {chunk['page_number']} - {chunk['section_title']}]:\n"
            f"{truncated_text}\n"
        )
    prompt += "\nFinal Answer:"
    return prompt

#########################
# Answer Generation via Ollama with Llama 3.2 3B
#########################
def generate_answer_with_ollama(prompt):
    """
    Run the Llama 3.2 3B model via the provided Ollama command.
    """
    try:
        result = subprocess.run(
            ["ollama", "run", "llama3.2:3b"],
            input=prompt.encode("utf-8"),
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=True
        )
        output_text = result.stdout.decode("utf-8")
        if "Final Answer:" in output_text:
            output_text = output_text.split("Final Answer:")[-1].strip()
        return output_text
    except subprocess.CalledProcessError as e:
        st.error("Error generating answer via Ollama:")
        st.error(e.stderr.decode("utf-8"))
        return None

#########################
# Initialize Chat History in Session State
#########################
if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = ""

#########################
# Streamlit UI
#########################
st.set_page_config(page_title="Document Query Assistant", layout="wide")
st.title("Document Query Assistant")
st.write("""
Upload one or more documents (PDF, TXT, DOCX, HTML, PPTX, or images) and then ask a question about their content.
The Llama 3.2 3B model will provide a final answer based solely on the documents.
""")

# Multi-file uploader to allow multiple document uploads
uploaded_files = st.file_uploader("Upload your documents", 
                                  type=["pdf", "txt", "docx", "html", "pptx", "png", "jpg", "jpeg", "tiff"], 
                                  accept_multiple_files=True)

all_pages_data = []
if uploaded_files:
    for uploaded_file in uploaded_files:
        st.info(f"Processing file: {uploaded_file.name}")
        pages_data = process_file(uploaded_file)
        all_pages_data.extend(pages_data)
    st.success(f"Processed {len(uploaded_files)} files successfully!")
    
    # Segment document(s) into chunks
    chunks = segment_document(all_pages_data)
    st.write(f"Total document segments: {len(chunks)}")
    
    # Compute embeddings and build FAISS index
    with st.spinner("Computing embeddings and building index..."):
        embeddings = compute_embeddings(chunks)
        index = build_faiss_index(embeddings)
    st.success("Embeddings computed and FAISS index built.")
    
    # Display chat history (conversational context)
    with st.expander("Chat History (previous queries and answers)"):
        st.write(st.session_state["chat_history"])
    
    # Input for new query
    query = st.text_input("Enter your question about the documents:")
    
    # Option to include conversation history in the prompt (append previous Q&A)
    include_history = st.checkbox("Include conversation history in prompt", value=True)
    
    if st.button("Get Answer") and query:
        st.info("Processing query and retrieving context...")
        top_k = 5  # Fixed number of context chunks
        retrieved_chunks = retrieve_relevant_chunks(query, index, chunks, top_k=top_k)
        
        if not retrieved_chunks:
            st.warning("No relevant excerpts found. Try rephrasing your query.")
        else:
            with st.expander("Show Retrieved Excerpts"):
                for chunk in retrieved_chunks:
                    st.markdown(f"**Document: {chunk['document_name']} | Page {chunk['page_number']} - {chunk['section_title']}**")
                    st.write(chunk['chunk_text'][:300] + "...")
            
            # Build conversation context from chat history if option enabled
            conversation_context = ""
            if include_history and st.session_state["chat_history"]:
                conversation_context = st.session_state["chat_history"]
            
            prompt = construct_prompt(query, retrieved_chunks, conversation_context=conversation_context)
            
            st.info("Just a moment...")
            answer = generate_answer_with_ollama(prompt)
            if answer:
                st.markdown("### Final Answer")
                st.write(answer)
                # Append new query and answer to chat history
                new_entry = f"**Q:** {query}\n**A:** {answer}\n\n"
                st.session_state["chat_history"] += new_entry
            else:
                st.error("Failed to generate an answer.")
else:
    st.info("Please upload one or more documents to get started.")